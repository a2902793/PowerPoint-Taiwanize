from pptx import Presentation

search_str = "old text"
repl_str = "123"

prs = Presentation("example.pptx")
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            if(shape.text.find(search_str))!=-1:
                text_frame = shape.text_frame
                cur_text = text_frame.paragraphs[0].runs[0].text
                new_text = cur_text.replace(str(search_str), str(repl_str))
                text_frame.paragraphs[0].runs[0].text = new_text
prs.save("example2.pptx")


    
