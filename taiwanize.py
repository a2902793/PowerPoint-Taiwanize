#!/usr/bin/env python
from pptx import Presentation
from langconv import Converter
from pathlib import Path
import argparse, sys, os

def Traditional2Simplified(sentence):  # 繁體轉簡體
    sentence = Converter('zh-hans').convert(str(sentence))
    return sentence


def Simplified2Traditional(sentence):  # 簡體轉繁體
    sentence = Converter('zh-hant').convert(str(sentence))
    return sentence

def PowerPoint(filepath):
    prs = Presentation(filepath)
    for slide in prs.slides:                                         # 每張投影片
        for shape in slide.shapes:                                   # 每張投影片/每個物件
            if shape.has_text_frame:                                 # 判斷每張投影片/每個物件是否包含文字框
                for paragraph in shape.text_frame.paragraphs:        # 每張投影片/每個物件/每個文字段落
                    for run in paragraph.runs:                       # 每張投影片/每個物件/每個文字段落/每個文字流
                        run.text = Simplified2Traditional(run.text)  # 將每張投影片/每個物件/每個文字段落/每個文字流的簡體中文翻成繁體中文
                        """
                        不可直接用 text_frame.text 的原因是因為「格式會跑掉」
                        從層級上來分，slide > shape > text_frame > paragraph > run
                                     投影片   物件      文字框      文字段落   文字流
                        
                        * 間距(line_spacing)是在 paragraph 中設定
                        * 字體(font)、大小(size)、粗斜體跟顏色等在最小層級的 run 中設定

                        若是直接使用層級較大的 text_frame.text 換詞語，等於 paragraph.line_spacing、run.font、run.size 等等的參數未設定，
                        這樣會導致原本的格式被刷掉。若要維持一樣的格式要指定只換掉 run.text 這樣其餘沒被你動到格式都會保留。
                        """
    prs.save("example_TW.pptx")

def main():
    parser = argparse.ArgumentParser(description='【台灣化（Taiwanized），一個將中國用詞用語翻成台灣化的 Python 腳本。】')
    parser.add_argument('filename', metavar='F', nargs='*', help='檔案名稱，範例："python taiwanized.py example.pptx"')
    parser.add_argument('-p', '--ppt', '--pptx', action='store_true', help='將此資料夾內的所有 PowerPoint 檔台灣化')
    if len(sys.argv) == 1:
        parser.print_help()
        sys.exit(1)
    args = parser.parse_args()
    
    if args.filename:
        for file in args.filename:
            if Path(file).suffix == '.pptx':
                PowerPoint(file)
            else:
                print("目前不支援檔案: " + '"' + file + '"')
    if args.ppt:
        for file in os.listdir("./"):
            if Path(file).suffix in (".ppt", ".pptx"):
                PowerPoint(file)


if __name__ == "__main__":
    main()
